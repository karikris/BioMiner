"""Aggregate-only, reproducible exploratory report for GBIF occurrence parquet files."""

from __future__ import annotations

import hashlib
import json
import shutil
import textwrap
import zipfile
from datetime import datetime
from pathlib import Path
from typing import Any

import duckdb
import matplotlib
import polars as pl
from pptx import Presentation
from pptx.util import Inches


matplotlib.use("Agg")
from matplotlib import pyplot as plt

_REQUIRED_COLUMNS = (
    "year",
    "month",
    "decimalLatitude",
    "decimalLongitude",
    "family",
    "genus",
    "species",
    "scientificName",
    "datasetName",
    "publisher",
    "institutionCode",
    "recordedBy",
    "basisOfRecord",
    "issue",
    "identifiedBy",
)

_SPILL_DIRECTORY_NAME = "duckdb_spill"
_SPILL_MAX_SIZE = "2GiB"
_PPTX_TIMESTAMP = (1980, 1, 1, 0, 0, 0)
_PPTX_PROPERTIES_TIMESTAMP = datetime(2000, 1, 1)
_CHART_MAX_BARS = 12
_CHART_LABEL_WRAP_WIDTH = 26
_CHART_MAX_LABEL_LINES = 2

ChartArtifact = tuple[str, str, Path, int, int]

_COMPLETENESS_FIELDS = (
    "year",
    "month",
    "decimalLatitude",
    "decimalLongitude",
    "family",
    "genus",
    "species",
    "scientificName",
    "datasetName",
    "recordedBy",
    "identifiedBy",
)


def build_ground_zero_eda_run(
    source_path: str | Path,
    output_dir: str | Path,
    source_manifest_path: str | Path | None = None,
    top_n: int = 20,
) -> dict[str, Any]:
    """Create a new, aggregate-only EDA report directory from an occurrence Parquet.

    The input is scanned by DuckDB and is never copied or mutated.  A pre-existing
    output directory is treated as an overwrite attempt and rejected.
    """
    source = Path(source_path)
    output = Path(output_dir)
    if output.exists():
        raise FileExistsError(f"refusing to overwrite existing output directory: {output}")
    if not source.is_file():
        raise FileNotFoundError(f"source parquet does not exist: {source}")
    if top_n < 1:
        raise ValueError("top_n must be at least 1")

    source_manifest = _read_source_manifest(source_manifest_path)
    connection = duckdb.connect()
    try:
        schema = _source_schema(connection, source)
        missing = [column for column in _REQUIRED_COLUMNS if column not in schema]
        if missing:
            raise ValueError("source parquet is missing required GBIF columns: " + ", ".join(missing))
    finally:
        connection.close()

    output.mkdir(parents=True)
    spill_directory = output / _SPILL_DIRECTORY_NAME
    connection = duckdb.connect()
    try:
        _configure_spill_directory(connection, spill_directory)
        work = _build_aggregate_work(connection, source, top_n)
    finally:
        connection.close()
        shutil.rmtree(spill_directory, ignore_errors=True)

    work_path = output / "eda_work.parquet"
    work.write_parquet(work_path)
    csv_paths = _write_csvs(work, output)
    chart_artifacts = _write_charts(work, output)
    deck_path = _write_deck(work, chart_artifacts, output)

    inventory = _artifact_inventory(output)
    manifest = {
        "report": "ground_zero_eda",
        "top_n": top_n,
        "source": {
            "path": str(source.resolve()),
            "sha256": _sha256(source),
            "schema": schema,
        },
        "source_manifest": source_manifest,
        "query_strategy": {
            "engine": "duckdb",
            "input": "read-only parquet_scan",
            "completeness": "single wide aggregate scan",
            "spill": {
                "temporary_directory": _SPILL_DIRECTORY_NAME,
                "max_size": _SPILL_MAX_SIZE,
                "cleaned": True,
            },
        },
        "chart_policy": {
            "max_bars": _CHART_MAX_BARS,
            "label_wrap_width": _CHART_LABEL_WRAP_WIDTH,
            "max_label_lines": _CHART_MAX_LABEL_LINES,
        },
        "chart_metadata": [
            {
                "path": path.name,
                "section": section,
                "metric": metric,
                "displayed_bars": displayed_bars,
                "available_ranked_values": available_ranked_values,
            }
            for section, metric, path, displayed_bars, available_ranked_values in chart_artifacts
        ],
        "artifact_inventory": inventory,
        "artifacts": {
            "work_parquet": work_path.name,
            "csv": [path.name for path in csv_paths],
            "charts": [path.name for _, _, path, _, _ in chart_artifacts],
            "deck": deck_path.name,
        },
    }
    # This is deliberately last: the manifest inventories every generated artifact.
    (output / "manifest.json").write_text(
        json.dumps(manifest, indent=2, sort_keys=True) + "\n", encoding="utf-8"
    )
    return {"output_dir": str(output), "manifest": manifest}


def _source_schema(connection: duckdb.DuckDBPyConnection, source: Path) -> dict[str, str]:
    rows = connection.execute("DESCRIBE SELECT * FROM parquet_scan(?)", [str(source)]).fetchall()
    return {str(name): _normalise_type(str(data_type)) for name, data_type, *_ in rows}


def _normalise_type(data_type: str) -> str:
    lowered = data_type.lower()
    if "char" in lowered or "string" in lowered or "text" in lowered:
        return "string"
    return lowered


def _configure_spill_directory(connection: duckdb.DuckDBPyConnection, spill_directory: Path) -> None:
    spill_directory.mkdir()
    escaped_directory = str(spill_directory.resolve()).replace("'", "''")
    connection.execute(f"SET temp_directory = '{escaped_directory}'")
    connection.execute(f"SET max_temp_directory_size = '{_SPILL_MAX_SIZE}'")


def _build_aggregate_work(
    connection: duckdb.DuckDBPyConnection, source: Path, top_n: int) -> pl.DataFrame:
    scans: list[pl.DataFrame] = [_completeness_query(connection, source)]
    scans.append(_issue_frequency_query(connection, source, top_n))

    for section, metric, field, null_label in (
        ("quality", "occurrences_by_basis_of_record", "basisOfRecord", "(NULL)"),
        ("temporal", "occurrences_by_year", "year", "(NULL)"),
        ("temporal", "occurrences_by_month", "month", "(NULL)"),
        ("taxonomy", "occurrences_by_family", "family", "(NULL)"),
        ("taxonomy", "occurrences_by_genus", "genus", "(NULL)"),
        ("taxonomy", "occurrences_by_species", "species", "(NULL)"),
        ("taxonomy", "occurrences_by_scientific_name", "scientificName", "(NULL)"),
        ("contributor", "occurrences_by_dataset", "datasetName", "(NULL)"),
        ("contributor", "occurrences_by_publisher", "publisher", "(NULL)"),
        ("contributor", "occurrences_by_institution", "institutionCode", "(NULL)"),
        ("contributor", "occurrences_by_recorder", "recordedBy", "(NULL)"),
    ):
        scans.append(
            _frequency_query(
                connection,
                source,
                section,
                metric,
                field,
                null_label,
                top_n,
                retain_null=(metric == "occurrences_by_month"),
            )
        )

    scans.append(_coordinate_quality_query(connection, source))
    scans.append(_geography_query(connection, source))
    work = pl.concat(scans, how="vertical")
    return work.sort(["section", "metric", "value", "dimension"], descending=[False, False, True, False])


def _completeness_query(connection: duckdb.DuckDBPyConnection, source: Path) -> pl.DataFrame:
    expressions = ", ".join(
        f"CAST(COUNT(*) FILTER (WHERE NULLIF(TRIM(CAST({_quote(field)} AS VARCHAR)), '') IS NOT NULL) AS BIGINT) AS {_quote(field)}"
        for field in _COMPLETENESS_FIELDS
    )
    row = connection.execute(
        f"SELECT CAST(COUNT(*) AS BIGINT) AS record_count, {expressions} FROM parquet_scan(?)",
        [str(source)],
    ).fetchone()
    assert row is not None
    total = int(row[0])
    return pl.DataFrame(
        {
            "section": ["completeness"] * len(_COMPLETENESS_FIELDS),
            "metric": ["nonempty_by_field"] * len(_COMPLETENESS_FIELDS),
            "dimension": list(_COMPLETENESS_FIELDS),
            "value": [int(value) for value in row[1:]],
            "record_count": [total] * len(_COMPLETENESS_FIELDS),
        }
    )


def _issue_frequency_query(
    connection: duckdb.DuckDBPyConnection, source: Path, top_n: int
) -> pl.DataFrame:
    sql = """
        WITH source AS (
            SELECT NULLIF(TRIM(CAST("issue" AS VARCHAR)), '') AS issues
            FROM parquet_scan(?)
        ), tokens AS (
            SELECT TRIM(token) AS dimension, CAST(COUNT(*) AS BIGINT) AS value
            FROM source, UNNEST(string_split(issues, ';')) AS split(token)
            WHERE TRIM(token) <> ''
            GROUP BY 1
        ), ranked AS (
            SELECT *, ROW_NUMBER() OVER (ORDER BY value DESC, dimension ASC) AS row_number
            FROM tokens
        ), null_bucket AS (
            SELECT '(NULL)' AS dimension, CAST(COUNT(*) AS BIGINT) AS value
            FROM source WHERE issues IS NULL
        ), selected AS (
            SELECT dimension, value FROM ranked WHERE row_number <= ?
            UNION ALL
            SELECT dimension, value FROM null_bucket
        )
        SELECT 'quality' AS section, 'occurrences_by_issue' AS metric, dimension, value,
               CAST((SELECT COUNT(*) FROM source) AS BIGINT) AS record_count
        FROM selected ORDER BY value DESC, dimension ASC
    """
    return pl.from_arrow(connection.execute(sql, [str(source), top_n]).to_arrow_table())


def _frequency_query(
    connection: duckdb.DuckDBPyConnection,
    source: Path,
    section: str,
    metric: str,
    field: str,
    null_label: str,
    top_n: int,
    retain_null: bool = False,
) -> pl.DataFrame:
    dimension = f"COALESCE(NULLIF(TRIM(CAST({_quote(field)} AS VARCHAR)), ''), '{null_label}')"
    sql = f"""
        WITH grouped AS (
            SELECT {dimension} AS dimension, CAST(COUNT(*) AS BIGINT) AS value
            FROM parquet_scan(?)
            GROUP BY 1
        ), ranked AS (
            SELECT *, ROW_NUMBER() OVER (ORDER BY value DESC, dimension ASC) AS row_number
            FROM grouped
        )
        SELECT '{section}' AS section, '{metric}' AS metric, dimension, value,
               CAST((SELECT SUM(value) FROM grouped) AS BIGINT) AS record_count
        FROM ranked WHERE row_number <= ? {"OR dimension = '(NULL)'" if retain_null else ""}
        ORDER BY value DESC, dimension ASC
    """
    return pl.from_arrow(connection.execute(sql, [str(source), top_n]).to_arrow_table())


def _coordinate_quality_query(connection: duckdb.DuckDBPyConnection, source: Path) -> pl.DataFrame:
    lat = "TRY_CAST(NULLIF(TRIM(CAST(\"decimalLatitude\" AS VARCHAR)), '') AS DOUBLE)"
    lon = "TRY_CAST(NULLIF(TRIM(CAST(\"decimalLongitude\" AS VARCHAR)), '') AS DOUBLE)"
    dimension = f"""CASE
        WHEN {lat} IS NULL OR {lon} IS NULL THEN 'missing_or_invalid'
        WHEN {lat} NOT BETWEEN -90 AND 90 OR {lon} NOT BETWEEN -180 AND 180 THEN 'out_of_range'
        ELSE 'valid'
    END"""
    return _grouped_expression_query(connection, source, "quality", "coordinate_quality", dimension)


def _geography_query(connection: duckdb.DuckDBPyConnection, source: Path) -> pl.DataFrame:
    lat = "TRY_CAST(NULLIF(TRIM(CAST(\"decimalLatitude\" AS VARCHAR)), '') AS DOUBLE)"
    dimension = f"""CASE
        WHEN {lat} IS NULL OR {lat} NOT BETWEEN -90 AND 90 THEN 'missing_or_invalid'
        WHEN {lat} < -30 THEN 'south_of_30S'
        WHEN {lat} < 0 THEN '30S_to_equator'
        WHEN {lat} < 30 THEN 'equator_to_30N'
        ELSE 'north_of_30N'
    END"""
    return _grouped_expression_query(connection, source, "geography", "occurrences_by_latitude_band", dimension)


def _grouped_expression_query(
    connection: duckdb.DuckDBPyConnection,
    source: Path,
    section: str,
    metric: str,
    dimension_sql: str,
) -> pl.DataFrame:
    sql = f"""
        WITH grouped AS (
            SELECT {dimension_sql} AS dimension, CAST(COUNT(*) AS BIGINT) AS value
            FROM parquet_scan(?) GROUP BY 1
        )
        SELECT '{section}' AS section, '{metric}' AS metric, dimension, value,
               CAST((SELECT SUM(value) FROM grouped) AS BIGINT) AS record_count
        FROM grouped ORDER BY value DESC, dimension ASC
    """
    return pl.from_arrow(connection.execute(sql, [str(source)]).to_arrow_table())


def _write_csvs(work: pl.DataFrame, output: Path) -> list[Path]:
    paths: list[Path] = []
    for section in sorted(work["section"].unique().to_list()):
        path = output / f"{section}_metrics.csv"
        section_frame = work.filter(pl.col("section") == section).with_columns(
            pl.col("dimension").map_elements(_csv_safe, return_dtype=pl.String)
        )
        path.write_bytes(b"\xef\xbb\xbf" + section_frame.write_csv().encode("utf-8"))
        paths.append(path)
    return paths


def _csv_safe(value: str | None) -> str | None:
    if value and value[0] in "=+-@":
        return "'" + value
    return value


def _write_charts(work: pl.DataFrame, output: Path) -> list[ChartArtifact]:
    """Write one readable chart for each material aggregate metric."""
    artifacts: list[ChartArtifact] = []
    metric_pairs = work.select("section", "metric").unique().sort(["section", "metric"])
    for section, metric in metric_pairs.iter_rows():
        full_frame = work.filter((pl.col("section") == section) & (pl.col("metric") == metric))
        available_ranked_values = full_frame.height
        frame = full_frame.head(_CHART_MAX_BARS)
        labels = [_chart_label(dimension) for dimension in frame["dimension"].to_list()]
        figure_width = 14.5
        figure_height = min(9.0, max(3.6, 3.3 + 0.47 * frame.height))
        longest_label = max((max(map(len, label.splitlines())) for label in labels), default=1)
        left_margin = min(0.42, max(0.30, 0.0145 * longest_label))
        fig, axis = plt.subplots(figsize=(figure_width, figure_height))
        positions = list(range(frame.height))
        axis.barh(positions, frame["value"].to_list(), color="#2f6f9f")
        axis.set_yticks(positions, labels)
        axis.invert_yaxis()
        axis.set_title(
            f"{section.title()} — {metric.replace('_', ' ').title()}\n"
            f"Displaying top {frame.height} of {available_ranked_values} ranked values",
            pad=16,
            fontsize=14,
        )
        axis.set_xlabel("Occurrences")
        axis.tick_params(axis="y", labelsize=8)
        axis.grid(axis="x", alpha=0.2)
        axis.set_axisbelow(True)
        fig.subplots_adjust(left=left_margin, right=0.97, top=0.84, bottom=0.11)
        path = output / f"{section}__{metric}.png"
        fig.savefig(path, dpi=150)
        plt.close(fig)
        artifacts.append((section, metric, path, frame.height, available_ranked_values))
    return artifacts


def _chart_label(dimension: str) -> str:
    """Wrap category labels deliberately so chart margins remain predictable."""
    lines = textwrap.wrap(
        str(dimension),
        width=_CHART_LABEL_WRAP_WIDTH,
        break_long_words=True,
        break_on_hyphens=False,
    ) or [""]
    if len(lines) > _CHART_MAX_LABEL_LINES:
        lines = lines[:_CHART_MAX_LABEL_LINES]
        suffix = str(dimension).strip()[-4:]
        available = _CHART_LABEL_WRAP_WIDTH - len(suffix) - 2
        lines[-1] = lines[-1][:available].rstrip() + "… " + suffix
    return "\n".join(lines)


def _write_deck(work: pl.DataFrame, charts: list[ChartArtifact], output: Path) -> Path:
    presentation = Presentation()
    _set_deterministic_core_properties(presentation)
    title = presentation.slides.add_slide(presentation.slide_layouts[0])
    title.shapes.title.text = "Ground Zero EDA"
    title.placeholders[1].text = f"Aggregate-only summary of {int(work['record_count'].max())} occurrences"
    for section, metric, chart, _, _ in charts:
        slide = presentation.slides.add_slide(presentation.slide_layouts[5])
        slide.shapes.title.text = f"{section.title()} — {metric.replace('_', ' ').title()}"
        slide.shapes.add_picture(str(chart), Inches(0.6), Inches(1.1), width=Inches(8.8))
    path = output / "ground_zero_eda.pptx"
    presentation.save(path)
    _canonicalize_pptx(path)
    return path


def _set_deterministic_core_properties(presentation: Presentation) -> None:
    properties = presentation.core_properties
    properties.author = "BioMiner"
    properties.category = "Ground Zero EDA"
    properties.comments = "Aggregate-only report"
    properties.created = _PPTX_PROPERTIES_TIMESTAMP
    properties.keywords = "BioMiner, GBIF, EDA"
    properties.last_modified_by = "BioMiner"
    properties.modified = _PPTX_PROPERTIES_TIMESTAMP
    properties.revision = 1
    properties.subject = "Aggregate occurrence exploratory data analysis"
    properties.title = "Ground Zero EDA"


def _canonicalize_pptx(path: Path) -> None:
    """Rewrite the OOXML ZIP with stable entry ordering and metadata."""
    canonical = path.with_suffix(".canonical.pptx")
    with zipfile.ZipFile(path, "r") as source, zipfile.ZipFile(
        canonical, "w", compression=zipfile.ZIP_DEFLATED, compresslevel=9
    ) as destination:
        for name in sorted(source.namelist()):
            info = zipfile.ZipInfo(filename=name, date_time=_PPTX_TIMESTAMP)
            info.compress_type = zipfile.ZIP_DEFLATED
            info.create_system = 3
            info.external_attr = 0o600 << 16
            destination.writestr(info, source.read(name), compress_type=zipfile.ZIP_DEFLATED, compresslevel=9)
    canonical.replace(path)


def _artifact_inventory(output: Path) -> list[dict[str, Any]]:
    return [
        {"path": path.relative_to(output).as_posix(), "bytes": path.stat().st_size, "sha256": _sha256(path)}
        for path in sorted(output.rglob("*"))
        if path.is_file() and path.name != "manifest.json"
    ]


def _sha256(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as handle:
        for chunk in iter(lambda: handle.read(1024 * 1024), b""):
            digest.update(chunk)
    return digest.hexdigest()


def _read_source_manifest(path: str | Path | None) -> dict[str, Any] | None:
    if path is None:
        return None
    manifest_path = Path(path)
    if not manifest_path.is_file():
        raise FileNotFoundError(f"source manifest does not exist: {manifest_path}")
    payload = json.loads(manifest_path.read_text(encoding="utf-8"))
    if not isinstance(payload, dict):
        raise ValueError("source manifest must be a JSON object")
    return payload


def _quote(identifier: str) -> str:
    return '"' + identifier.replace('"', '""') + '"'


__all__ = ["build_ground_zero_eda_run"]
